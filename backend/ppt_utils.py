from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def generate_pptx_deck(slides_data: list[dict], output_path: str | Path) -> None:
    """Generates a professional 16:9 widescreen presentation from slide bullet details.

    Each dict in slides_data should contain:
    - "title": Title string
    - "bullets": List of bullet points (strings)
    """
    prs = Presentation()

    # Set 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Cohesive clean palette
    color_bg = RGBColor(248, 250, 252)       # Slate 50 (soft light gray)
    color_primary = RGBColor(15, 23, 42)      # Slate 900 (deep charcoal)
    color_secondary = RGBColor(71, 85, 105)   # Slate 600 (cool gray)
    color_accent = RGBColor(16, 163, 127)     # Assistant Emerald (accent green)

    # Use layout 6 (blank slide) to custom-position everything precisely
    blank_layout = prs.slide_layouts[6]

    for index, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)

        # Apply solid background fill
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_bg

        is_title_slide = (index == 0)

        if is_title_slide:
            # 1. Title Slide Layout
            # Title Text Box
            title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.2))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            tf_title.margin_left = 0
            tf_title.margin_top = 0

            p_title = tf_title.paragraphs[0]
            p_title.text = slide_info.get("title", "Research Presentation").upper()
            p_title.font.name = "Arial"
            p_title.font.size = Pt(40)
            p_title.font.bold = True
            p_title.font.color.rgb = color_primary

            # Subtitle Box
            sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.333), Inches(2.0))
            tf_sub = sub_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = 0

            bullets = slide_info.get("bullets", [])
            for idx, bullet in enumerate(bullets):
                p_sub = tf_sub.add_paragraph() if idx > 0 else tf_sub.paragraphs[0]
                p_sub.text = bullet
                p_sub.font.name = "Arial"
                p_sub.font.size = Pt(18)
                # First line is emerald accent, rest is cool gray
                p_sub.font.color.rgb = color_accent if idx == 0 else color_secondary
                p_sub.space_after = Pt(8)
        else:
            # 2. Content Slide Layout
            # Slide Header Title
            header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
            tf_header = header_box.text_frame
            tf_header.word_wrap = True
            tf_header.margin_left = 0
            tf_header.margin_top = 0

            p_header = tf_header.paragraphs[0]
            p_header.text = slide_info.get("title", "Slide Title")
            p_header.font.name = "Arial"
            p_header.font.size = Pt(32)
            p_header.font.bold = True
            p_header.font.color.rgb = color_primary

            # Content Text Box
            content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.6))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            tf_content.margin_left = 0

            bullets = slide_info.get("bullets", [])
            for idx, bullet in enumerate(bullets):
                p_bullet = tf_content.add_paragraph() if idx > 0 else tf_content.paragraphs[0]
                p_bullet.text = bullet
                p_bullet.font.name = "Arial"
                p_bullet.font.size = Pt(16)
                p_bullet.font.color.rgb = color_primary
                p_bullet.space_after = Pt(14)
                p_bullet.level = 0

            # Slide Footer
            footer_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.8), Inches(11.333), Inches(0.4))
            tf_footer = footer_box.text_frame
            tf_footer.margin_right = 0
            p_foot = tf_footer.paragraphs[0]
            p_foot.text = f"Multimodal AI Research Assistant  |  Slide {index + 1}"
            p_foot.font.name = "Arial"
            p_foot.font.size = Pt(9)
            p_foot.font.color.rgb = color_secondary
            p_foot.alignment = PP_ALIGN.RIGHT

    prs.save(str(output_path))
